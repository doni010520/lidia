"""Processador de mídia — converte mídia recebida em texto para o pipeline.

TODAS as funções usam uaz.download_message(message_id) para obter bytes
decriptografados. Mídia do WhatsApp vem criptografada (.enc) — acessar
a URL direta falha. O endpoint /message/download da uazapi faz o decrypt.

Resposta de /message/download contém campo 'file' com bytes em base64.
"""
from __future__ import annotations

import base64
import io

import openai
from loguru import logger

from app.core.config import settings
from app.services.deps import get_uaz_client


async def _download_media_bytes(message_id: str) -> bytes | None:
    """Baixa mídia via uazapi /message/download → retorna bytes decriptografados."""
    uaz = get_uaz_client()
    try:
        result = await uaz.download_message(message_id)
        file_b64 = result.get("file", "")
        if not file_b64:
            logger.warning(f"download_message sem campo 'file' para {message_id}")
            return None
        return base64.b64decode(file_b64)
    except Exception:
        logger.exception(f"Erro ao baixar mídia {message_id}")
        return None


async def transcribe_audio(message_id: str) -> str:
    """Transcreve áudio via uazapi (Whisper integrado)."""
    uaz = get_uaz_client()
    try:
        result = await uaz.download_message(
            message_id,
            transcribe=True,
            openai_apikey=settings.openai_api_key,
        )
        transcription = result.get("transcription", "")
        if transcription:
            logger.debug(f"Áudio transcrito ({len(transcription)} chars)")
        return transcription
    except Exception:
        logger.exception("Erro na transcrição de áudio")
        return ""


async def describe_image(message_id: str, mimetype: str | None = None) -> str:
    """Descreve imagem via GPT Vision usando base64 data URL."""
    content = await _download_media_bytes(message_id)
    if content is None:
        return "[Imagem recebida, mas não foi possível baixá-la]"

    mime = mimetype or "image/jpeg"
    b64 = base64.b64encode(content).decode("utf-8")
    data_url = f"data:{mime};base64,{b64}"

    client = openai.AsyncOpenAI(api_key=settings.openai_api_key)
    try:
        response = await client.chat.completions.create(
            model=settings.openai_model,
            max_tokens=500,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": (
                                "Descreva esta imagem de forma objetiva e concisa em português. "
                                "Se contiver texto, transcreva-o. "
                                "Se for um documento, extraia as informações principais."
                            ),
                        },
                        {
                            "type": "image_url",
                            "image_url": {"url": data_url, "detail": "auto"},
                        },
                    ],
                }
            ],
        )
        description = response.choices[0].message.content or ""
        logger.debug(f"Imagem descrita ({len(description)} chars)")
        return f"[Imagem recebida: {description}]"
    except Exception:
        logger.exception("Erro ao descrever imagem via GPT Vision")
        return "[Imagem recebida, mas não foi possível processá-la]"


async def extract_document(message_id: str, mimetype: str | None = None) -> str:
    """Extrai texto de documento (PDF, DOCX, XLSX, PPTX).

    Baixa bytes decriptografados via uazapi e processa com lib adequada.
    """
    content = await _download_media_bytes(message_id)
    if content is None:
        return "[Documento recebido, mas não foi possível baixá-lo]"

    mime = (mimetype or "").lower()

    # PDF
    if "pdf" in mime:
        try:
            import fitz
            doc = fitz.open(stream=content, filetype="pdf")
            pages = [page.get_text() for page in doc]
            doc.close()
            text = "\n\n".join(pages).strip()
            if text:
                return f"[Documento PDF recebido]\n{text[:3000]}"
            return "[Documento PDF recebido, mas sem texto extraível]"
        except Exception:
            logger.exception("Erro ao processar PDF")
            return "[Documento PDF recebido, mas houve erro no processamento]"

    # DOCX
    if "wordprocessing" in mime or "docx" in mime:
        try:
            from docx import Document
            doc = Document(io.BytesIO(content))
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            if text:
                return f"[Documento Word recebido]\n{text[:3000]}"
            return "[Documento Word recebido, mas sem texto]"
        except Exception:
            logger.exception("Erro ao processar DOCX")
            return "[Documento Word recebido, mas houve erro no processamento]"

    # XLSX
    if "spreadsheet" in mime or "xlsx" in mime:
        try:
            from openpyxl import load_workbook
            wb = load_workbook(io.BytesIO(content), read_only=True)
            lines = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    line = " | ".join(str(c) for c in row if c is not None)
                    if line.strip():
                        lines.append(line)
            wb.close()
            text = "\n".join(lines[:100])
            if text:
                return f"[Planilha recebida]\n{text[:3000]}"
            return "[Planilha recebida, mas sem dados]"
        except Exception:
            logger.exception("Erro ao processar XLSX")
            return "[Planilha recebida, mas houve erro no processamento]"

    # PPTX
    if "presentation" in mime or "pptx" in mime:
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        texts.append(shape.text_frame.text)
            text = "\n".join(texts)
            if text:
                return f"[Apresentação recebida]\n{text[:3000]}"
            return "[Apresentação recebida, mas sem texto]"
        except Exception:
            logger.exception("Erro ao processar PPTX")
            return "[Apresentação recebida, mas houve erro no processamento]"

    return f"[Documento recebido (tipo: {mime}), formato não suportado para extração]"


async def handle_video(message_id: str) -> str:
    """Download de vídeo via uazapi → upload para Google Drive → retorna file_id."""
    from app.services import drive_client

    content = await _download_media_bytes(message_id)
    if content is None:
        return ""

    try:
        file_id = drive_client.upload_bytes(
            content,
            filename=f"video_{message_id}.mp4",
            mimetype="video/mp4",
            folder_id=settings.drive_folder_media,
        )
        logger.info(f"Vídeo uploaded para Drive: {file_id}")
        return file_id
    except Exception:
        logger.exception("Erro ao fazer upload do vídeo para Drive")
        return ""
