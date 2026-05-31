"""Indexar base de conhecimento no pgvector.

Uso:
    # Indexar documentos locais
    python -m scripts.index_knowledge --path ./docs/

    # Indexar de pasta do Drive (requer GOOGLE_SERVICE_ACCOUNT_JSON)
    python -m scripts.index_knowledge --folder <drive_folder_id>

    # Limpar base antes de indexar
    python -m scripts.index_knowledge --path ./docs/ --clear
"""
from __future__ import annotations

import argparse
import asyncio
import json
import os
from pathlib import Path

import openai
from loguru import logger
from sqlalchemy import text

from app.core.config import settings
from app.db import engine

# ── Configuração de chunking ──
CHUNK_SIZE = 1000
CHUNK_OVERLAP = 200


def chunk_text(content: str, source: str, *, size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[dict]:
    """Divide texto em chunks com overlap."""
    chunks = []
    start = 0
    idx = 0
    while start < len(content):
        end = start + size
        chunk = content[start:end].strip()
        if chunk:
            chunks.append({
                "content": chunk,
                "source": source,
                "metadata": {"source": source, "chunk_index": idx},
            })
            idx += 1
        start = end - overlap
    return chunks


def extract_text_from_file(filepath: Path) -> str:
    """Extrai texto de arquivo (PDF, DOCX, XLSX, PPTX, TXT)."""
    suffix = filepath.suffix.lower()

    if suffix == ".txt" or suffix == ".md":
        return filepath.read_text(encoding="utf-8", errors="replace")

    if suffix == ".pdf":
        try:
            import fitz  # pymupdf
            doc = fitz.open(str(filepath))
            pages = [page.get_text() for page in doc]
            doc.close()
            return "\n\n".join(pages)
        except ImportError:
            logger.warning("pymupdf não instalado, pulando PDF")
            return ""

    if suffix == ".docx":
        try:
            from docx import Document
            doc = Document(str(filepath))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except ImportError:
            logger.warning("python-docx não instalado, pulando DOCX")
            return ""

    if suffix == ".xlsx":
        try:
            from openpyxl import load_workbook
            wb = load_workbook(str(filepath), read_only=True)
            lines = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    line = " | ".join(str(c) for c in row if c is not None)
                    if line.strip():
                        lines.append(line)
            wb.close()
            return "\n".join(lines)
        except ImportError:
            logger.warning("openpyxl não instalado, pulando XLSX")
            return ""

    if suffix == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(str(filepath))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        texts.append(shape.text_frame.text)
            return "\n".join(texts)
        except ImportError:
            logger.warning("python-pptx não instalado, pulando PPTX")
            return ""

    logger.warning(f"Formato não suportado: {suffix}")
    return ""


async def embed_texts(texts: list[str]) -> list[list[float]]:
    """Gera embeddings em batch."""
    client = openai.AsyncOpenAI(api_key=settings.openai_api_key)
    embeddings = []

    # Batch de 100 (limite da API)
    for i in range(0, len(texts), 100):
        batch = texts[i:i + 100]
        resp = await client.embeddings.create(
            model=settings.openai_embedding_model,
            input=batch,
        )
        embeddings.extend([d.embedding for d in resp.data])

    return embeddings


async def index_local(path: str, clear: bool = False) -> None:
    """Indexa documentos de diretório local."""
    folder = Path(path)
    if not folder.exists():
        logger.error(f"Caminho não existe: {path}")
        return

    supported = {".txt", ".md", ".pdf", ".docx", ".xlsx", ".pptx"}
    files = [f for f in folder.rglob("*") if f.suffix.lower() in supported]

    if not files:
        logger.warning(f"Nenhum arquivo suportado em {path}")
        return

    logger.info(f"Encontrados {len(files)} arquivos para indexar")

    # Extrair e chunkar
    all_chunks: list[dict] = []
    for filepath in files:
        logger.info(f"Processando: {filepath.name}")
        content = extract_text_from_file(filepath)
        if not content.strip():
            logger.warning(f"  → sem conteúdo, pulando")
            continue
        chunks = chunk_text(content, filepath.name)
        all_chunks.extend(chunks)
        logger.info(f"  → {len(chunks)} chunks")

    if not all_chunks:
        logger.warning("Nenhum chunk gerado")
        return

    logger.info(f"Total: {len(all_chunks)} chunks. Gerando embeddings...")

    # Gerar embeddings
    texts = [c["content"] for c in all_chunks]
    embeddings = await embed_texts(texts)

    # Inserir no banco
    async with engine.begin() as conn:
        if clear:
            await conn.execute(text("DELETE FROM knowledge_chunks"))
            logger.info("Base limpa (--clear)")

        for chunk, emb in zip(all_chunks, embeddings):
            await conn.execute(
                text("""
                    INSERT INTO knowledge_chunks (content, embedding, metadata, source)
                    VALUES (:content, :embedding::vector, :metadata::jsonb, :source)
                """),
                {
                    "content": chunk["content"],
                    "embedding": str(emb),
                    "metadata": json.dumps(chunk["metadata"]),
                    "source": chunk["source"],
                },
            )

    logger.info(f"✅ {len(all_chunks)} chunks indexados com sucesso")
    await engine.dispose()


async def main() -> None:
    parser = argparse.ArgumentParser(description="Indexar base de conhecimento")
    parser.add_argument("--path", help="Caminho local com documentos")
    parser.add_argument("--folder", help="ID da pasta no Google Drive")
    parser.add_argument("--clear", action="store_true", help="Limpar base antes de indexar")
    args = parser.parse_args()

    if args.path:
        await index_local(args.path, clear=args.clear)
    elif args.folder:
        logger.info("Indexação via Drive será implementada junto com drive_client (Fase 5)")
    else:
        parser.print_help()


if __name__ == "__main__":
    asyncio.run(main())
